import re
from datetime import datetime
import pandas as pd
from pptx.dml.color import RGBColor
from pathlib import Path

def load_excel(file_path, sheet_name=0):
    """
    Load an Excel file into a DataFrame.

    Parameters:
    file_path (str): The path to the Excel file.
    sheet_name (str or int, optional): The sheet name or index to load. Defaults to the first sheet.

    Returns:
    DataFrame: The loaded DataFrame.
    """
    return pd.read_excel(file_path, sheet_name=sheet_name)


def save_excel(df, file_path):
    """
    Save a DataFrame to an Excel file.

    Parameters:
    df (DataFrame): The DataFrame to save.
    file_path (str): The path to the Excel file.
    """
    df.to_excel(file_path, index=False)


def process_dataframe(df):
    """
    Process the DataFrame by cleaning, renaming columns, removing NaN, restructuring,
    converting data types, and formating.

    Parameters:
    df (DataFrame): The DataFrame to process.

    Returns:
    DataFrame: The processed DataFrame.
    """
    
    header_row = locate_header_row(df)
    df = clean_dataframe(df, header_row)
    df = restructure_dataframe(df)
    processed_df = format_dataframe(df)
    return processed_df


def locate_header_row(df):
    """
    Locate the header row in the DataFrame dynamically.

    Parameters:
    df (DataFrame): The DataFrame to search.

    Returns:
    int: The index of the header row.
    """
    return df[df.apply(
        lambda row: row.astype(str).str.contains(r'Điểm dự báo', flags=re.IGNORECASE).any(), axis=1)].index[0]


def clean_dataframe(df, header_row):
    """
    Clean the DataFrame by setting the header row and removing unnecessary rows and columns.

    Parameters:
    df (DataFrame): The DataFrame to clean.
    header_row (int): The index of the header row.

    Returns:
    DataFrame: The cleaned DataFrame.
    """
    
    # Locate and remove whitespaces in header row
    df.columns = df.iloc[header_row]
    df = df.drop(index=range(0, header_row+1))
    df.columns = df.columns.str.strip()
    
    # Find columns with 'NaN' as header
    nan_indices = [i for i, col in enumerate(df.columns) if pd.isna(col)]

    # Keep only the first 'NaN' column
    if len(nan_indices) > 1:
        keep_indices = [i for i in range(df.shape[1]) if i not in nan_indices[1:]]
        df = df.iloc[:, keep_indices]

    # Clean column names
    mapping = {
        "Thứ hai": "Thứ 2",
        "Thứ ba": "Thứ 3",
        "Thứ tư": "Thứ 4",
        "Thứ năm": "Thứ 5",
        "Thứ sáu": "Thứ 6",
        "Thứ bảy": "Thứ 7"
    }
    df = df.rename(columns=mapping)

    headers = df.columns
    date_row = df.iloc[0]
    date_row = date_row.apply(
        lambda x: datetime.strptime(str(x), "%Y-%m-%d %H:%M:%S").strftime("%d/%m/%Y") if pd.notnull(x) else pd.NaT
    )
    headers = [
        f'Ngày {value} ({header})' if pd.notnull(value) else header
        for header, value in zip(headers, date_row)
    ]
    df.columns = headers
    df = df.drop(index=header_row+1)

    # Remove rows that are all NaN
    df = df.replace(r'^\s*$', pd.NA, regex=True)
    df = df.dropna(how='all')

    # Remove the signature row
    df = df.iloc[:-1]

    # Fill NaN values in the 1st and 2nd column using ffill 
    df.iloc[:, 0] = df.iloc[:, 0].ffill()
    df.iloc[:, 1] = df.iloc[:, 1].ffill()

    # If the values from the second column are not NaN, concatenate them to the 3nd column and separate value by '_'. 
    # After that , remove the 3rd column

    df.iloc[:, 1] = df.apply(
        lambda row: f'{row.iloc[1]} {row.iloc[2]}' if pd.notna(row.iloc[2]) else row.iloc[1], axis=1)
    df = df.drop(df.columns[2], axis=1)

    # Remove any remaining rows with NaN
    df = df.dropna(how='any')

    # Normalize the column headers by make sure only the first letter is capitalized
    df.columns = [col.capitalize() for col in df.columns]

    # Remove special characters in the second column, which will be used as headers later
    df.iloc[:, 1] = df.iloc[:, 1].apply(lambda x: re.sub(r'\(°C\)|%|\(mm\)', '', str(x)))
    df.iloc[:, 1] = df.iloc[:, 1].apply(
        lambda x: ' '.join(
            [x.split()[0].capitalize()] + [word.lower() for word in x.split()[1:]]))

    return df


def restructure_dataframe(df):
    """
    Restructure the DataFrame by grouping and transposing.

    Parameters:
    df (DataFrame): The DataFrame to restructure.

    Returns:
    DataFrame: The restructured DataFrame.
    """
    grouped = df.groupby('Điểm dự báo')
    sub_dfs = []
    for name, group in grouped:
        sub_df = group.drop(columns=['Điểm dự báo']).reset_index(drop=True)
        sub_df = sub_df.transpose()
        sub_df.columns = sub_df.iloc[0]
        sub_df = sub_df.drop(sub_df.index[0])
        sub_df.insert(0, 'Ngày', sub_df.index)
        sub_df.insert(0, 'Điểm dự báo', name)
        sub_dfs.append(sub_df)
    return pd.concat(sub_dfs, ignore_index=False)


def format_dataframe(df):
    """
    Format the DataFrame for output.

    Parameters:
    df (DataFrame): The DataFrame to format.

    Returns:
    DataFrame: The formatted DataFrame.
    """
    output = df[
        ['Điểm dự báo', 'Ngày', 'Thời tiết', 'Nhiệt độ cao nhất', 
         'Nhiệt độ thấp nhất']].copy()
    output.iloc[:, 1] = output.iloc[:, 1] + '\n' + output.iloc[:, 2]
    output['Nhiệt độ cao nhất'] = output['Nhiệt độ cao nhất'].astype(str)
    output['Nhiệt độ thấp nhất'] = output['Nhiệt độ thấp nhất'].astype(str)
    output['Nhiệt độ'] = output.iloc[:, 3] + '°C\n' + output.iloc[:, 4] + '°C'
    output = output.drop(output.columns[2:5], axis=1)
    return output


def find_table(slide):
    """
    Find the first table in the slide.

    Parameters:
    slide (Slide): The slide to search.

    Returns:
    Table: The first table found in the slide, or None if no table is found.
    """
    for shape in slide.shapes:
        if shape.has_table:
            return shape.table
    return None


def get_run_styles(paragraph):
    """
    Extract styles from paragraph runs.

    Parameters:
    paragraph (Paragraph): The paragraph to extract styles from.

    Returns:
    list: A list of dictionaries containing the styles of each run.
    """
    return [
        {
            'text': run.text,
            'size': run.font.size,
            'bold': run.font.bold,
            'italic': run.font.italic,
            'underline': run.font.underline,
            'name': run.font.name,
            'color': run.font.color.rgb if run.font.color and hasattr(run.font.color, 'rgb') else None
        }
        for run in paragraph.runs
    ]


def apply_styles(paragraph, text, styles, temp=True):
    """
    Apply styles to paragraph text.

    Parameters:
    paragraph (Paragraph): The paragraph to apply styles to.
    text (str): The text to apply styles to.
    styles (list): A list of dictionaries containing the styles to apply.
    temp (bool, optional): Whether to apply temporary styles. Defaults to True.
    """
    paragraph.clear()
    words = text.split()
    for i, word in enumerate(words):
        style = styles[i % len(styles)]
        new_run = paragraph.add_run()
        new_run.text = word if i == len(words) - 1 else word + ' '
        new_run.font.size = style['size']
        new_run.font.bold = style['bold']
        new_run.font.italic = style['italic']
        new_run.font.underline = style['underline']
        new_run.font.name = style['name']
        if style['color']:
            new_run.font.color.rgb = style['color']
        elif temp:
            new_run.font.color.rgb = RGBColor(0x44, 0x72, 0xC4)


def set_cell_text_two_paragraphs(cell, new_text, temp=True):
    """
    Set cell text with two paragraphs, preserving styles.

    Parameters:
    cell (Cell): The cell to set text in.
    new_text (str): The new text to set.
    temp (bool, optional): Whether to apply temporary styles. Defaults to True.
    """
    parts = new_text.split('\n')
    text1 = parts[0] if len(parts) > 0 else ''
    text2 = parts[1] if len(parts) > 1 else ''

    while len(cell.text_frame.paragraphs) < 2:
        cell.text_frame.add_paragraph()

    p1, p2 = cell.text_frame.paragraphs[:2]

    if p1.runs:
        styles = get_run_styles(p1)
        apply_styles(p1, text1, styles, temp)
    else:
        p1.text = text1

    if p2.runs:
        styles = get_run_styles(p2)
        apply_styles(p2, text2, styles, temp)
    else:
        p2.text = text2


def update_table_with_data(table, df):
    """
    Update table cells with data from DataFrame.

    Parameters:
    table (Table): The table to update.
    df (DataFrame): The DataFrame containing the data.
    """

    for i in range(10):
        date = re.sub(r'/\d{4}', '', df.iloc[i, 1])
        set_cell_text_two_paragraphs(table.cell(i, 0), date, temp=False)
        set_cell_text_two_paragraphs(table.cell(i, 2), df.iloc[i, 2])


def get_row_heights(table):
    """
    Get row heights of the table.

    Parameters:
    table (Table): The table to get row heights from.

    Returns:
    list: A list of row heights.
    """
    tbl = table._graphic_frame._element
    tr_elements = tbl.xpath('.//a:tr')
    return [int(tr.get('h') or 0) for tr in tr_elements]


def remove_all_pictures(slide):
    """
    Remove all pictures from the slide.

    Parameters:
    slide (Slide): The slide to remove pictures from.
    """
    for shape in list(slide.shapes):
        if shape.shape_type == 13:
            slide.shapes._spTree.remove(shape._element)


def calculate_cell_position(table, row, col):
    """
    Calculate the position of a cell in the table.

    Parameters:
    table (Table): The table containing the cell.
    row (int): The row index of the cell.
    col (int): The column index of the cell.

    Returns:
    tuple: The (left, top, width, height) position of the cell.
    """
    table_frame = table._graphic_frame
    table_x = table_frame._element.xfrm.off.x
    table_y = table_frame._element.xfrm.off.y

    cell_left = table_x
    cell_top = table_y

    for r in range(row):
        cell_top += table.rows[r].height

    for c in range(col):
        cell_left += table.columns[c].width

    cell_height = table.rows[row].height
    cell_width = table.columns[col].width

    return cell_left + (cell_width - cell_width * 0.8) / 2, \
           cell_top + (cell_height - cell_height * 0.8) / 2, \
           cell_width * 0.8, \
           cell_height * 0.8


def add_image_to_cell(slide, table, row, col, image_path):
    """
    Add an image to a specific cell in the table.

    Parameters:
    slide (Slide): The slide containing the table.
    table (Table): The table containing the cell.
    row (int): The row index of the cell.
    col (int): The column index of the cell.
    image_path (str): The path to the image file.
    """
    left, top, width, height = calculate_cell_position(table, row, col)
    slide.shapes.add_picture(image_path, left, top, width, height)


def extract_conditions(text, phrases):
    """
    Extract weather conditions from text.

    Parameters:
    text (str): The text to extract conditions from.
    phrases (list): A list of phrases to extract.

    Returns:
    set: A set of extracted conditions.
    """
    extracted = set()

    for phrase in phrases:
        if phrase in text:
            extracted.add(phrase)
            text = text.replace(phrase, "")

    extracted.update(re.findall(r'\b\w+\b', text.lower()))
    return extracted


def select_image(text, phrases, image_mappings):
    """
    Select an image based on weather conditions.

    Parameters:
    text (str): The text containing weather conditions.
    phrases (list): A list of phrases to extract conditions.
    image_mappings (dict): A dictionary mapping conditions to image filenames.

    Returns:
    str: The selected image filename, or None if no match is found.
    """
    text_set = extract_conditions(text, phrases)

    for condition, image in image_mappings.items():
        if condition.issubset(text_set):
            return image

    return None


def update_table_with_images(slide, table, phrases, image_mappings, image_dir):
    """
    Update table cells with images based on conditions.

    Parameters:
    slide (Slide): The slide containing the table.
    table (Table): The table to update.
    phrases (list): A list of phrases to extract conditions.
    image_mappings (dict): A dictionary mapping conditions to image filenames.
    image_dir (Path): The directory containing the images.
    """
    for i in range(10):
        img = select_image(table.cell(i, 0).text, phrases, image_mappings)
        if img:
            img_path = str(Path(image_dir) / img)
            add_image_to_cell(slide, table, i, 1, img_path)


def write_district(slide, new_text):
    """
    Write district name in the slide.

    Parameters:
    slide (Slide): The slide to write the district name in.
    new_text (str): The district name to write.

    Returns:
    bool: True if the district name was written, False otherwise.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape.text = new_text
            return True
    return False


def extract_period(df):
    """
    Extract the start and end dates from the second column of the DataFrame.

    The start date is extracted from the first row, and the end date is extracted from the tenth row.
    Both dates are in the format 'Ngày dd/mm', where 'dd' is the day and 'mm' is the month.

    If the start and end dates are in the same month, the function returns:
        'Từ ngày {start_day} - {end_day} tháng {month} năm 2025'
    Otherwise, it returns:
        'Từ ngày {start_day}/{start_month} - {end_day}/{end_month} năm 2025'

    Parameters:
    df (DataFrame): The DataFrame to extract dates from.

    Returns:
    str: The formatted period string.
    """
    start_text = df.iloc[0, 1]
    end_text = df.iloc[9, 1]

    start_date = re.search(r'Ngày (\d{2}/\d{2}/\d{4})', start_text).group(1)
    end_date = re.search(r'Ngày (\d{2}/\d{2}/\d{4})', end_text).group(1)

    start_day, start_month, start_year = start_date.split('/')
    end_day, end_month, end_year = end_date.split('/')

    start_day = start_day.lstrip('0')
    end_day = end_day.lstrip('0')
    start_month = start_month.lstrip('0')
    end_month = end_month.lstrip('0')

    if start_month == end_month:
        return f'Từ ngày {start_day} - {end_day} tháng {start_month} năm {end_year}'
    else:
        return f'Từ ngày {start_day}/{start_month} - {end_day}/{end_month} năm {end_year}'


def write_period(slide, new_text):
    """
    Write period in the slide.

    Parameters:
    slide (Slide): The slide to write the period in.
    new_text (str): The period text to write.
    """
    p = slide.shapes[2].text_frame.paragraphs[0]
    styles = get_run_styles(p)
    apply_styles(p, new_text, styles, temp=False)


def extract_province(filename):
    """
    Extract the province substring from the filename before the first underscore.

    Parameters:
    filename (str): The filename to extract the substring from.

    Returns:
    str: The extracted substring, or None if no match is found.
    """
    match = re.match(r'([^_]+)_', filename)
    return match.group(1) if match else None